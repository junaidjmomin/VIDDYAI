"""
Content Generation Router
PowerPoint generation and YouTube video search
"""

from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
from pptx import Presentation
from services.llm_service import generate_key_points
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import tempfile
import os
from dotenv import load_dotenv
import uuid
import base64
import requests

# YouTube API
try:
    from googleapiclient.discovery import build
    YOUTUBE_AVAILABLE = True
except ImportError:
    YOUTUBE_AVAILABLE = False
    print("Warning: google-api-python-client not installed. YouTube search will use fallback.")

load_dotenv()
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

if not OPENAI_API_KEY:
    raise ValueError("OPENAI_API_KEY not found in environment variables")

router = APIRouter()


class PPTRequest(BaseModel):
    concept: str
    grade: int
    subject: Optional[str] = "Science"
    include_activity: bool = True


@router.post("/generate/ppt")
async def generate_ppt(data: PPTRequest):
    """
    Generate a PowerPoint presentation for a concept
    
    Creates a beautiful, grade-appropriate PPT with:
    - Title slide
    - Key points slides (one per point)
    - Activity slide
    
    Returns: PPT file download
    """

    key_points = generate_key_points(
    concept=data.concept,
    grade=data.grade,
    subject=data.subject
)
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Color theme - VidyaSetu brand colors
    bg_dark = RGBColor(0x0A, 0x0E, 0x1A)
    bg_medium = RGBColor(0x11, 0x18, 0x27)
    accent_yellow = RGBColor(0xFC, 0xD3, 0x4D)
    accent_purple = RGBColor(0x8B, 0x5C, 0xF6)
    accent_cyan = RGBColor(0x06, 0xB6, 0xD4)
    accent_green = RGBColor(0x10, 0xB9, 0x81)
    accent_orange = RGBColor(0xF5, 0x9E, 0x0B)
    text_light = RGBColor(0xF1, 0xF5, 0xF9)
    text_muted = RGBColor(0x94, 0xA3, 0xB8)
    
    # SLIDE 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = bg_dark
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    tf = title_box.text_frame
    tf.text = data.concept
    tf.paragraphs[0].runs[0].font.size = Pt(60)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = accent_yellow
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(0.8))
    stf = subtitle_box.text_frame
    stf.text = f"Grade {data.grade} | {data.subject}"
    stf.paragraphs[0].runs[0].font.size = Pt(28)
    stf.paragraphs[0].runs[0].font.color.rgb = text_muted
    stf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # VidyaSetu branding
    brand_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.5))
    btf = brand_box.text_frame
    btf.text = "VidyaSetu AI 🦉"
    btf.paragraphs[0].runs[0].font.size = Pt(20)
    btf.paragraphs[0].runs[0].font.color.rgb = accent_purple
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # SLIDES 2-N: Key Points
    colors = [accent_purple, accent_cyan, accent_green, accent_orange]
    
    for i, point in enumerate(key_points[:6]):  # Max 6 key points
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = bg_medium
        
        # Step number badge
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1.2), Inches(1))
        ntf = num_box.text_frame
        ntf.text = f"{i+1:02d}"
        ntf.paragraphs[0].runs[0].font.size = Pt(48)
        ntf.paragraphs[0].runs[0].font.bold = True
        color = colors[i % len(colors)]
        ntf.paragraphs[0].runs[0].font.color.rgb = color
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Key point text
        pt_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(10), Inches(3))
        ptf = pt_box.text_frame
        ptf.word_wrap = True
        ptf.text = point
        ptf.paragraphs[0].runs[0].font.size = Pt(40) if data.grade <= 3 else Pt(36)
        ptf.paragraphs[0].runs[0].font.bold = True
        ptf.paragraphs[0].runs[0].font.color.rgb = text_light
        
        # Decorative line
        line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(2), Inches(5.8), Inches(10), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.line.color.rgb = color
    
    # LAST SLIDE: Activity (if requested)
    if data.include_activity:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = bg_dark
        
        # Activity icon/emoji
        emoji_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2.33), Inches(1))
        etf = emoji_box.text_frame
        etf.text = "✏️"
        etf.paragraphs[0].runs[0].font.size = Pt(72)
        etf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Activity prompt
        act_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10.33), Inches(2.5))
        atf = act_box.text_frame
        atf.text = f"Try It Yourself!\n\nCan you explain '{data.concept}' to a friend?"
        atf.paragraphs[0].runs[0].font.size = Pt(44)
        atf.paragraphs[0].runs[0].font.bold = True
        atf.paragraphs[0].runs[0].font.color.rgb = accent_yellow
        atf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Bonus prompt
        bonus_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9.33), Inches(1))
        botf = bonus_box.text_frame
        botf.text = "Bonus: Draw a picture of this concept in your notebook! 🎨"
        botf.paragraphs[0].runs[0].font.size = Pt(24)
        botf.paragraphs[0].runs[0].font.color.rgb = text_muted
        botf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save to temporary file
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp_file.name)
    tmp_file.close()
    
    # Clean filename
    safe_filename = "".join(c if c.isalnum() or c in (' ', '-', '_') else '_' for c in data.concept)
    filename = f"{safe_filename}_Grade{data.grade}.pptx"
    
    return FileResponse(
        tmp_file.name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
        headers={
            "Content-Disposition": f"attachment; filename={filename}"
        }
    )


class ImageRequest(BaseModel):
    concept: str
    grade: int
    subject: str

@router.post("/generate/images")
async def generate_images(data: ImageRequest):

    images = []

    for i in range(3):
        prompt = f"""
                    Create a clear, colorful educational diagram explaining {data.concept}
                    for a Grade {data.grade} {data.subject} student.
                    Style: simple textbook illustration, labeled parts, kid-friendly,
                    white background, no watermark, high clarity.
                    """

        image_url = generate_image_file(prompt)
        images.append(image_url)

    return {
        "success": True,
        "images": images
    } 


def generate_image_file(prompt: str):
    import uuid
    import requests
    import os

    os.makedirs("static/generated", exist_ok=True)

    encoded_prompt = requests.utils.quote(prompt)

    image_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=1024&height=1024&seed=42"

    headers = {
        "User-Agent": "Mozilla/5.0"
    }

    response = requests.get(image_url, headers=headers, timeout=60)

    print("STATUS:", response.status_code)

    if response.status_code != 200:
        print("RESPONSE:", response.text)
        raise Exception("Free image generation failed")

    file_name = f"{uuid.uuid4()}.png"
    file_path = f"static/generated/{file_name}"

    with open(file_path, "wb") as f:
        f.write(response.content)

    return f"http://localhost:8000/static/generated/{file_name}"

def build_search_query(concept, grade, subject, iq_score, eq_score):
    base = f"CBSE Grade {grade} {subject} {concept}"

    if iq_score < 40:
        style = "simple explanation for kids step by step"
    elif iq_score > 75:
        style = "advanced detailed explanation with experiments"
    else:
        style = "clear explanation with examples"

    if eq_score < 40:
        tone = "real life examples storytelling"
    elif eq_score > 75:
        tone = "interactive discussion animation"
    else:
        tone = "engaging animated explanation"

    return f"{base} {style} {tone}"

@router.get("/video/search")
async def search_video(
    concept: str,
    grade: int,
    subject: str,
    iq_score: float,
    eq_score: float
):
    """
    Search for educational videos on YouTube
    
    Uses YouTube Data API v3 to find kid-safe educational videos
    
    Query params:
        concept: Topic to search for
        grade: Student's grade (1-5)
        subject: Subject area (default: Science)
        
    Returns:
        Video ID, embed URL, title, thumbnail
    """
    
    if not YOUTUBE_AVAILABLE:
        # Fallback response if YouTube API not available
        return {
            "success": False,
            "video_id": None,
            "embed_url": None,
            "title": "Video search unavailable",
            "message": "YouTube API client not installed"
        }
    
    youtube_api_key = os.getenv("YOUTUBE_API_KEY")
    
    if not youtube_api_key:
        # Fallback if API key not configured
        return {
            "success": False,
            "video_id": None,
            "embed_url": None,
            "title": "Video search unavailable",
            "message": "YouTube API key not configured. Add YOUTUBE_API_KEY to .env"
        }
    
    try:
        # Build YouTube API client
        youtube = build("youtube", "v3", developerKey=youtube_api_key)
        
        # Construct search query optimized for kids
        search_query = build_search_query(concept, grade, subject, iq_score, eq_score)
        
        # Search for videos
        search_response = youtube.search().list(
            q=search_query,
            part="snippet",
            maxResults=3,  # Get top 3 to pick best one
            type="video",
            videoDuration="short",  # Prefer short videos (< 4 min)
            safeSearch="strict",  # Only kid-safe content
            relevanceLanguage="en",
            order="relevance"
        ).execute()
        
        if not search_response.get("items"):
            # No results found
            return {
                "success": False,
                "video_id": None,
                "embed_url": None,
                "title": f"No videos found for '{concept}'",
                "message": "Try a different search term"
            }
        
        # Pick first result
        video = search_response["items"][0]
        video_id = video["id"]["videoId"]
        title = video["snippet"]["title"]
        thumbnail = video["snippet"]["thumbnails"]["medium"]["url"]
        
        return {
            "success": True,
            "video_id": video_id,
            "embed_url": f"https://www.youtube.com/embed/{video_id}?autoplay=0&modestbranding=1",
            "watch_url": f"https://www.youtube.com/watch?v={video_id}",
            "title": title,
            "thumbnail": thumbnail,
            "description": video["snippet"]["description"][:200] + "..."
        }
        
    except Exception as e:
        print(f"YouTube API error: {e}")
        
        # Fallback response
        return {
            "success": False,
            "video_id": None,
            "embed_url": None,
            "title": "Video search temporarily unavailable",
            "message": str(e),
            "error": True
        }


@router.get("/video/embed/{video_id}")
async def get_video_embed(video_id: str):
    """
    Get embed URL and metadata for a specific YouTube video
    """
    return {
        "success": True,
        "video_id": video_id,
        "embed_url": f"https://www.youtube.com/embed/{video_id}?autoplay=0&modestbranding=1",
        "watch_url": f"https://www.youtube.com/watch?v={video_id}"
    }
